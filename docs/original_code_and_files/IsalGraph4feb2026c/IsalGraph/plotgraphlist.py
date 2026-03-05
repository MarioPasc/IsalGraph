from typing import List, Tuple, Dict
from math import cos, sin, pi

# From package python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

from sparsegraph import SparseGraph
from circulardoublylinkedlist import CircularDoublyLinkedList


def _circular_layout(
    node_count: int,
    center: Tuple[float, float],
    radius: float
) -> List[Tuple[float, float]]:
    """
    Compute a circular layout for graph visualization.

    Nodes are placed evenly on a circle.

    :param node_count:
        Number of nodes in the graph.
    :param center:
        (x, y) center of the circle in inches.
    :param radius:
        Radius of the circle in inches.

    :return:
        A list of (x, y) positions in inches, indexed by node ID.
    """
    # Unpack center coordinates
    cx, cy = center

    # Allocate list for node positions
    positions: List[Tuple[float, float]] = []

    # Compute position for each node
    for i in range(node_count):
        # Evenly spaced angle around the circle
        angle: float = 2.0 * pi * i / node_count

        # Convert polar coordinates to Cartesian
        x: float = cx + radius * cos(angle)
        y: float = cy + radius * sin(angle)

        # Store the position
        positions.append((x, y))

    return positions


def old_graphs_to_pptx(
    graphs: List[SparseGraph],
    output_path: str
) -> None:
    """
    Generate a PowerPoint slideshow with one slide per SparseGraph.

    Each slide contains:
    - A title identifying the graph
    - Nodes drawn as labeled circles
    - Edges drawn as straight-line connectors

    A simple circular layout is used for deterministic and fast rendering.

    :param graphs:
        A list of SparseGraph instances to visualize.
    :param output_path:
        File path where the PowerPoint (.pptx) will be saved.
    """
    # Create a new PowerPoint presentation
    presentation: Presentation = Presentation()

    # Select a blank slide layout (no placeholders)
    blank_layout = presentation.slide_layouts[6]

    # Iterate over all graphs
    for graph_index, graph in enumerate(graphs):
        # Add a new slide
        slide = presentation.slides.add_slide(blank_layout)

        # ------------------------------------------------------------------
        # Add slide title
        # ------------------------------------------------------------------

        # Create a text box for the title
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.2),
            Inches(9.0),
            Inches(0.6)
        )

        # Access the text frame
        title_frame = title_box.text_frame
        title_frame.clear()

        # Set title text
        title_frame.text = f"Graph {graph_index + 1}"

        # Format title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(28)
        title_paragraph.alignment = PP_ALIGN.CENTER

        # ------------------------------------------------------------------
        # Compute node layout
        # ------------------------------------------------------------------

        # Get the number of nodes in the graph
        node_count: int = graph.node_count()

        # Compute circular layout positions
        positions = _circular_layout(
            node_count=node_count,
            center=(5.0, 4.0),   # slide center in inches
            radius=2.5
        )

        # ------------------------------------------------------------------
        # Draw edges
        # ------------------------------------------------------------------

        # Iterate over all source nodes
        for source in range(node_count):
            # Get the source node position
            x1, y1 = positions[source]

            # Iterate over outgoing edges
            for target in graph.neighbors(source):
                # Get the target node position
                x2, y2 = positions[target]

                # Draw a straight connector between nodes
                slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1),
                    Inches(y1),
                    Inches(x2),
                    Inches(y2)
                )

        # ------------------------------------------------------------------
        # Draw nodes
        # ------------------------------------------------------------------

        # Radius of node circle in inches
        node_radius: float = 0.2

        # Draw each node
        for node_id, (x, y) in enumerate(positions):
            # Create a circular shape for the node
            node_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - node_radius),
                Inches(y - node_radius),
                Inches(2 * node_radius),
                Inches(2 * node_radius)
            )

            # Add node label
            node_shape.text_frame.text = str(node_id)
            node_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save the PowerPoint file
    presentation.save(output_path)


def trace_to_pptx(
    trace: List[Tuple[SparseGraph, CircularDoublyLinkedList, int, int]],
    output_path: str
) -> None:
    """
    Generate a PowerPoint slideshow with one slide per SparseGraph.

    In addition to plotting the graph structure, this function:
    - Draws links induced by a CircularDoublyLinkedList as purple dashed edges
    - Highlights specific graph nodes with colored squares

    :param graphs:
        List of SparseGraph instances to visualize.
    :param cdll:
        CircularDoublyLinkedList whose node data reference graph node indices.
    :param start_index:
        Index of a CDLL node whose referenced graph node is highlighted in green.
    :param end_index:
        Index of a CDLL node whose referenced graph node is highlighted in red.
    :param output_path:
        File path where the PowerPoint file will be saved.
    """
    presentation: Presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    for graph_index, (graph, cdll, primary_ptr, secondary_ptr, current_str) in enumerate(trace):
        slide = presentation.slides.add_slide(blank_layout)

        # ------------------------------------------------------------
        # Title
        # ------------------------------------------------------------
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
        )
        title_box.text_frame.text = f"Graph {graph_index + 1}, {current_str}"
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # ------------------------------------------------------------
        # Layout
        # ------------------------------------------------------------
        node_count = graph.node_count()
        center = (5.0, 4.0)
        radius = 2.5
        node_radius = 0.2

        positions: List[Tuple[float, float]] = []
        for i in range(node_count):
            angle = 2 * pi * i / node_count
            positions.append((
                center[0] + radius * cos(angle),
                center[1] + radius * sin(angle)
            ))

        # ------------------------------------------------------------
        # Determine highlighted nodes
        # ------------------------------------------------------------
        green_node = cdll.get_value(primary_ptr)
        blue_node = cdll.get_value(cdll.next_node(primary_ptr))
        red_node = cdll.get_value(secondary_ptr)

        highlight_colors: Dict[int, RGBColor] = {
            green_node: RGBColor(0, 176, 80),
            blue_node: RGBColor(0, 112, 192),
            red_node: RGBColor(192, 0, 0),
        }

        # ------------------------------------------------------------
        # Create grouped nodes
        # ------------------------------------------------------------
        node_circles: Dict[int, object] = {}

        for node_id, (x, y) in enumerate(positions):
            group = slide.shapes.add_group_shape()

            # Optional square (only if highlighted)
            if node_id in highlight_colors:
                square = group.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x - 0.3),
                    Inches(y - 0.3),
                    Inches(0.6),
                    Inches(0.6),
                )
                square.fill.background()
                square.line.color.rgb = highlight_colors[node_id]
                square.line.width = Pt(2)

            # Circle (always)
            circle = group.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - node_radius),
                Inches(y - node_radius),
                Inches(2 * node_radius),
                Inches(2 * node_radius),
            )
            circle.text_frame.text = str(node_id)
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Store circle reference for connectors
            node_circles[node_id] = circle

        # ------------------------------------------------------------
        # Graph edges
        # ------------------------------------------------------------
        for u in range(node_count):
            for v in graph.neighbors(u):
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0
                )
                conn.begin_connect(node_circles[u], 0)
                conn.end_connect(node_circles[v], 0)

        # ------------------------------------------------------------
        # CDLL links (purple dashed)
        # ------------------------------------------------------------
        for i in range(cdll.size()):
            u = cdll.get_value(i)
            v = cdll.get_value(cdll.next_node(i))

            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0
            )
            conn.begin_connect(node_circles[u], 0)
            conn.end_connect(node_circles[v], 0)
            conn.line.color.rgb = RGBColor(128, 0, 128)
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    presentation.save(output_path)